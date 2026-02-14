"""Advanced content summarization with batch processing and intelligent model routing.

This module extends the basic SummaryEngine with:
- Grouped batch processing by file type (process all text, then all images, etc.)
- Intelligent model loading/unloading to manage VRAM
- Support for multiple file types: images, PDFs, Office docs, videos, audio
- Fallback strategies when primary methods fail
"""

from __future__ import annotations

import base64
import tempfile
from dataclasses import dataclass
from enum import Enum
from pathlib import Path

from duplicleaner.ai.summaries import SummaryEngine
from duplicleaner.db.database import Database
from duplicleaner.db.models import FileRecord
from duplicleaner.utils.config import get_config
from duplicleaner.utils.lmstudio_manager import LMStudioManager, ModelType
from duplicleaner.utils.logging import get_logger

logger = get_logger(__name__)


class ProcessingType(Enum):
    """File processing categories for batch grouping."""
    TEXT = "text"           # Plain text files -> text model
    IMAGE = "image"         # Photos, images -> vision model
    VISUAL_DOC = "visual"   # PDFs, Office docs -> convert to images -> vision model
    VIDEO = "video"         # Videos -> extract keyframes -> vision model
    AUDIO = "audio"         # Audio -> transcribe -> text model
    SKIP = "skip"           # Unsupported file types


# File extension mappings
TEXT_EXTENSIONS = {
    '.txt', '.md', '.rst', '.log', '.csv', '.tsv',
    '.py', '.js', '.java', '.cpp', '.c', '.h', '.cs', '.go', '.rs', '.php',
    '.json', '.xml', '.yaml', '.yml', '.toml', '.ini', '.conf', '.cfg',
    '.html', '.css', '.sql', '.sh', '.bat', '.ps1',
    '.eml', '.msg',  # Email (if plain text extraction works)
}

IMAGE_EXTENSIONS = {
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif',
    '.heic', '.heif', '.ico', '.svg',
}

VISUAL_DOC_EXTENSIONS = {
    '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt',
    '.odt', '.ods', '.odp',  # LibreOffice
}

VIDEO_EXTENSIONS = {
    '.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv', '.m4v',
    '.mpeg', '.mpg', '.3gp',
}

AUDIO_EXTENSIONS = {
    '.mp3', '.wav', '.m4a', '.flac', '.ogg', '.aac', '.wma', '.opus',
}

# Files we intentionally skip (not useful to summarize)
SKIP_EXTENSIONS = {
    '.exe', '.dll', '.so', '.dylib', '.bin',  # Binaries
    '.zip', '.rar', '.7z', '.tar', '.gz',      # Archives
    '.db', '.sqlite', '.mdb',                  # Databases
    '.iso', '.dmg', '.img',                    # Disk images
    '.tmp', '.cache', '.lock',                 # Temp files
}


@dataclass
class BatchProgress:
    """Progress tracking for batch summarization."""
    total_files: int = 0
    processed_files: int = 0
    current_phase: str = "initializing"
    current_file: str = ""
    text_files: int = 0
    image_files: int = 0
    visual_doc_files: int = 0
    video_files: int = 0
    audio_files: int = 0
    skipped_files: int = 0
    successful: int = 0
    failed: int = 0

    @property
    def percent_complete(self) -> float:
        if self.total_files == 0:
            return 0.0
        return (self.processed_files / self.total_files) * 100


class ContentSummarizer:
    """Advanced summarizer with intelligent batch processing and model management."""

    def __init__(self, db: Database) -> None:
        self.db = db
        self.config = get_config()
        self.engine = SummaryEngine(db)
        self.progress = BatchProgress()
        self._whisper_model = None
        self._faster_whisper_model = None

        # Track loaded model type to avoid unnecessary reloads
        self._current_model_type: ProcessingType | None = None

        # LMStudio manager for automatic model detection/switching
        self.lmstudio_manager: LMStudioManager | None = None
        if self.config.ai.summary_provider == "lmstudio":
            try:
                self.lmstudio_manager = LMStudioManager()
                if not self.lmstudio_manager.is_available():
                    logger.warning("LMStudioMonitorService not available - model detection disabled")
                    self.lmstudio_manager = None
                else:
                    logger.info("LMStudioMonitorService connected - automatic model detection enabled")
            except Exception as exc:
                logger.warning("Failed to initialize LMStudioManager: %s", exc)
                self.lmstudio_manager = None

    def summarize_directory_batch(
        self,
        directory: str,
        file_types: set[str] | None = None,
        limit: int | None = None,
    ) -> BatchProgress:
        """
        Summarize all files in a directory using intelligent batch processing.

        Groups files by type and processes them in order:
        1. Text files (fastest, text model)
        2. Images (vision model)
        3. Visual documents (convert + vision model)
        4. Videos (keyframes + vision model)
        5. Audio (transcribe + text model)

        Args:
            directory: Directory path to process
            file_types: Optional set of extensions to filter (e.g., {'.jpg', '.pdf'})
            limit: Optional limit on total files to process

        Returns:
            BatchProgress with statistics
        """
        logger.info("Starting batch summarization for directory: %s", directory)

        # Get all file records from database for this directory
        file_records = self._get_files_for_directory(directory, file_types, limit)

        if not file_records:
            logger.warning("No files found in directory: %s", directory)
            return self.progress

        # Group files by processing type
        grouped = self._group_files_by_type(file_records)

        # Initialize progress tracking
        self.progress = BatchProgress(
            total_files=len(file_records),
            text_files=len(grouped[ProcessingType.TEXT]),
            image_files=len(grouped[ProcessingType.IMAGE]),
            visual_doc_files=len(grouped[ProcessingType.VISUAL_DOC]),
            video_files=len(grouped[ProcessingType.VIDEO]),
            audio_files=len(grouped[ProcessingType.AUDIO]),
            skipped_files=len(grouped[ProcessingType.SKIP]),
        )

        logger.info("File type breakdown: text=%d, image=%d, visual_doc=%d, video=%d, audio=%d, skip=%d",
                    self.progress.text_files, self.progress.image_files,
                    self.progress.visual_doc_files, self.progress.video_files,
                    self.progress.audio_files, self.progress.skipped_files)

        # Process each group in optimal order
        self._process_batch(ProcessingType.TEXT, grouped[ProcessingType.TEXT])
        self._process_batch(ProcessingType.IMAGE, grouped[ProcessingType.IMAGE])
        self._process_batch(ProcessingType.VISUAL_DOC, grouped[ProcessingType.VISUAL_DOC])
        self._process_batch(ProcessingType.VIDEO, grouped[ProcessingType.VIDEO])
        self._process_batch(ProcessingType.AUDIO, grouped[ProcessingType.AUDIO])

        self.progress.current_phase = "complete"
        logger.info("Batch summarization complete: %d successful, %d failed, %d skipped",
                    self.progress.successful, self.progress.failed, self.progress.skipped_files)

        return self.progress

    def _get_files_for_directory(
        self,
        directory: str,
        file_types: set[str] | None,
        limit: int | None,
    ) -> list[FileRecord]:
        """Query database for files in directory that need summaries."""
        # Get files from database that match directory and don't have summaries yet
        all_files = self.db.get_files_needing_summary_in_directory(
            directory,
            limit=999999,  # Get all files first, then we'll filter
            file_types=list(file_types) if file_types else None
        )

        # Filter by file types if specified
        if file_types:
            all_files = [f for f in all_files if f.file_type and f.file_type.lower() in file_types]

        # Filter out files that already have summaries
        files_needing_summaries = []
        for file_record in all_files:
            if file_record.id:
                existing = self.db.get_ai_summary(file_record.id)
                if not existing:
                    files_needing_summaries.append(file_record)

        # Apply limit
        if limit and len(files_needing_summaries) > limit:
            files_needing_summaries = files_needing_summaries[:limit]

        return files_needing_summaries

    def _group_files_by_type(self, file_records: list[FileRecord]) -> dict[ProcessingType, list[FileRecord]]:
        """Group files by their processing type."""
        groups: dict[ProcessingType, list[FileRecord]] = {
            ProcessingType.TEXT: [],
            ProcessingType.IMAGE: [],
            ProcessingType.VISUAL_DOC: [],
            ProcessingType.VIDEO: [],
            ProcessingType.AUDIO: [],
            ProcessingType.SKIP: [],
        }

        for record in file_records:
            ext = (record.file_type or "").lower()
            proc_type = self._classify_file_type(ext)
            groups[proc_type].append(record)

        return groups

    def _classify_file_type(self, extension: str) -> ProcessingType:
        """Determine processing type for a file extension."""
        if extension in TEXT_EXTENSIONS:
            return ProcessingType.TEXT
        if extension in IMAGE_EXTENSIONS:
            return ProcessingType.IMAGE
        if extension in VISUAL_DOC_EXTENSIONS:
            return ProcessingType.VISUAL_DOC
        if extension in VIDEO_EXTENSIONS:
            return ProcessingType.VIDEO
        if extension in AUDIO_EXTENSIONS:
            return ProcessingType.AUDIO
        if extension in SKIP_EXTENSIONS:
            return ProcessingType.SKIP

        # Unknown extension - check if it's text-like or skip
        if extension.startswith('.'):
            # Assume unknown text formats might be readable
            return ProcessingType.TEXT

        return ProcessingType.SKIP

    def _process_batch(self, proc_type: ProcessingType, file_records: list[FileRecord]) -> None:
        """Process a batch of files of the same type."""
        if not file_records:
            return

        if proc_type == ProcessingType.SKIP:
            self.progress.processed_files += len(file_records)
            return

        self.progress.current_phase = f"processing_{proc_type.value}_files"
        logger.info("Processing %d %s files", len(file_records), proc_type.value)

        # Load appropriate model for this batch
        self._ensure_model_loaded(proc_type)

        # Process each file
        for record in file_records:
            self.progress.current_file = record.path
            self.progress.processed_files += 1

            logger.info("[%d/%d] Processing %s: %s",
                        self.progress.processed_files, self.progress.total_files,
                        proc_type.value, record.path)

            try:
                success = self._process_single_file(proc_type, record)
                if success:
                    self.progress.successful += 1
                else:
                    logger.warning("[FAILED] %s: %s", proc_type.value, record.path)
                    self.progress.failed += 1
            except Exception as exc:
                logger.warning("[FAILED] %s: %s -- %s", proc_type.value, record.path, exc)
                self.progress.failed += 1

    def _ensure_model_loaded(self, proc_type: ProcessingType) -> None:
        """
        Ensure the correct model is loaded for the processing type.

        If LMStudioMonitorService is available, this will:
        - Detect if LMStudio has the wrong model loaded
        - Provide clear guidance on which model to load
        - Optionally wait for the user to switch models

        If LMStudioMonitorService is not available, provides manual instructions.
        """
        if self._current_model_type == proc_type:
            return  # Already loaded

        required_model_type = self._get_required_model_type(proc_type)

        # If LMStudioManager is available, check and guide
        if self.lmstudio_manager and required_model_type != ModelType.UNKNOWN:
            current_model = self.lmstudio_manager.get_current_model()

            if current_model:
                logger.info("Current model in LMStudio: %s (%s)", current_model.name, current_model.type.value)

                if current_model.type == required_model_type:
                    logger.info("Correct model type already loaded")
                    self._current_model_type = proc_type
                    return

                # Wrong model type loaded
                logger.warning(
                    "Wrong model type loaded! Current: %s, Required: %s",
                    current_model.type.value,
                    required_model_type.value,
                )
                logger.warning("=" * 60)
                logger.warning("PLEASE SWITCH MODEL IN LMSTUDIO:")
                logger.warning(self.lmstudio_manager.get_model_recommendation(required_model_type))
                logger.warning("=" * 60)
                logger.warning("Waiting 10 seconds for model switch...")

                # Give user time to switch
                import time
                time.sleep(10)

                # Check again
                new_model = self.lmstudio_manager.get_current_model()
                if new_model and new_model.type == required_model_type:
                    logger.info("Model switched successfully: %s", new_model.name)
                else:
                    logger.error("Model still not correct - processing may fail!")

            else:
                logger.warning("No model loaded in LMStudio!")
                logger.warning(self.lmstudio_manager.get_model_recommendation(required_model_type))

        else:
            # Fallback to manual instructions
            if proc_type in (ProcessingType.TEXT, ProcessingType.AUDIO):
                logger.info("=" * 60)
                logger.info("ENSURE TEXT MODEL IS LOADED IN LMSTUDIO")
                logger.info("Recommended: Josiefied-DeepSeek-R1-Qwen3-8B-abliterated")
                logger.info("Alternatives: Llama-3.2-3B-Instruct, Qwen2.5-3B-Instruct")
                logger.info("=" * 60)
            elif proc_type in (ProcessingType.IMAGE, ProcessingType.VISUAL_DOC, ProcessingType.VIDEO):
                logger.info("=" * 60)
                logger.info("ENSURE VISION MODEL IS LOADED IN LMSTUDIO")
                logger.info("Recommended: Qwen2.5-VL-7B")
                logger.info("Alternatives: LLaVA-v1.6-7B")
                logger.info("=" * 60)

        self._current_model_type = proc_type

    def _get_required_model_type(self, proc_type: ProcessingType) -> ModelType:
        """Map ProcessingType to ModelType."""
        if proc_type in (ProcessingType.TEXT, ProcessingType.AUDIO):
            return ModelType.TEXT
        elif proc_type in (ProcessingType.IMAGE, ProcessingType.VISUAL_DOC, ProcessingType.VIDEO):
            return ModelType.VISION
        else:
            return ModelType.UNKNOWN

    def _process_single_file(self, proc_type: ProcessingType, record: FileRecord) -> bool:
        """Process a single file based on its type."""
        try:
            if proc_type == ProcessingType.TEXT:
                return self._process_text_file(record)
            elif proc_type == ProcessingType.IMAGE:
                return self._process_image_file(record)
            elif proc_type == ProcessingType.VISUAL_DOC:
                return self._process_visual_doc(record)
            elif proc_type == ProcessingType.VIDEO:
                return self._process_video_file(record)
            elif proc_type == ProcessingType.AUDIO:
                return self._process_audio_file(record)
            return False
        except Exception as exc:
            logger.warning("Error processing %s: %s", record.path, exc)
            return False

    def _process_text_file(self, record: FileRecord) -> bool:
        """Process plain text file using text model."""
        try:
            # Read file content
            with open(record.path, encoding='utf-8', errors='ignore') as f:
                text = f.read()

            if not text.strip():
                logger.debug("Skipping empty file: %s", record.path)
                return False

            # Generate summary using text model
            summary = self.engine._generate_text_summary(text, file_path=record.path)
            if not summary:
                return False

            # Store in database
            from datetime import datetime

            from duplicleaner.db.models import AISummary

            ai_summary = AISummary(
                file_id=record.id,
                summary=summary,
                document_summary=summary,  # Same for text files
                summary_model=self.config.ai.get_summary_model(),
                generated_at=datetime.now(),
                user_edited=False,
            )
            self.db.add_ai_summary(ai_summary)
            return True

        except Exception as exc:
            logger.warning("Text processing failed for %s: %s", record.path, exc)
            return False

    def _process_image_file(self, record: FileRecord) -> bool:
        """Process image file using vision model."""
        summary = self.engine.analyze_file(record)
        if summary and summary.summary:
            logger.info("[OK] %s => %s", record.path, summary.summary[:120])
        return summary is not None

    def _process_visual_doc(self, record: FileRecord) -> bool:
        """
        Process visual document (PDF, Office) by converting to images.

        Strategy:
        1. Try text extraction first (faster for text-based PDFs)
        2. If text extraction fails or produces garbage, convert to images
        3. Process images with vision model
        """
        ext = (record.file_type or "").lower()

        if ext == '.pdf':
            return self._process_pdf(record)
        elif ext in {'.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt'}:
            return self._process_office_doc(record)

        return False

    def _process_pdf(self, record: FileRecord) -> bool:
        """Process PDF - try text extraction first, fall back to image rendering.

        When extracted JPEG pages exist in the database, uses those directly
        instead of rendering to temp files.
        """
        try:
            # Try text extraction first (faster)
            text = self._extract_pdf_text(record.path)

            if text and len(text.strip()) > 100 and self._is_readable_text(text):
                # Text-based PDF - use text model
                summary = self.engine._generate_text_summary(text, file_path=record.path)
                if summary:
                    from datetime import datetime

                    from duplicleaner.db.models import AISummary

                    ai_summary = AISummary(
                        file_id=record.id,
                        summary=summary,
                        document_summary=summary,
                        document_type="pdf",
                        summary_model=self.config.ai.get_summary_model(),
                        generated_at=datetime.now(),
                        user_edited=False,
                    )
                    self.db.add_ai_summary(ai_summary)
                    return True

            # Scanned PDF or text extraction failed - convert to images
            logger.info("PDF text extraction failed, converting to images: %s", record.path)

            # Check for pre-extracted JPEG pages
            image_paths = self._get_extracted_pdf_pages(record)
            if not image_paths:
                # Fall back to temp rendering
                with tempfile.TemporaryDirectory(prefix="duplicleaner_pdf_") as tmpdir:
                    image_paths = self._render_pdf_pages(record.path, tmpdir, max_pages=3)

            page_summaries = []
            for idx, image_path in enumerate(image_paths[:3], 1):
                summary = self.engine._generate_summary(image_path, original_path=record.path)
                if summary:
                    page_summaries.append(f"Page {idx}: {summary}")

            if not page_summaries:
                return False

            combined = self._combine_page_summaries(page_summaries, label="pdf", file_path=record.path)
            if not combined:
                combined = " ".join(s.split(": ", 1)[1] if ": " in s else s for s in page_summaries)

            from datetime import datetime

            from duplicleaner.db.models import AISummary

            ai_summary = AISummary(
                file_id=record.id,
                summary=combined,
                document_summary=combined,
                document_type="pdf",
                summary_model=self.config.ai.get_summary_model(),
                generated_at=datetime.now(),
                user_edited=False,
            )
            self.db.add_ai_summary(ai_summary)
            return True

        except Exception as exc:
            logger.warning("PDF processing failed for %s: %s", record.path, exc)
            return False

    def _get_extracted_pdf_pages(self, record: FileRecord) -> list[str]:
        """Get paths to pre-extracted JPEG pages for a PDF, if available.

        Returns:
            List of image file paths, or empty list if no extractions exist.
        """
        if not record.id:
            return []
        try:
            extractions = self.db.get_pdf_extractions(record.id)
            if not extractions:
                return []
            paths = []
            for _page_num, extracted_file_id in extractions:
                extracted_rec = self.db.get_file(extracted_file_id)
                if extracted_rec and Path(extracted_rec.path).exists():
                    paths.append(extracted_rec.path)
                else:
                    return []  # Incomplete extraction, fall back to rendering
            return paths
        except Exception:
            return []

    def _extract_pdf_text(self, pdf_path: str) -> str | None:
        """Extract text from PDF using pymupdf."""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(pdf_path)
            text_parts = []

            # Extract text from first 10 pages (limit for performance)
            for page_num in range(min(10, len(doc))):
                page = doc[page_num]
                text_parts.append(page.get_text())

            doc.close()
            return "\n".join(text_parts)

        except ImportError:
            logger.warning("PyMuPDF not installed - cannot extract PDF text")
            return None
        except Exception as exc:
            logger.warning("PDF text extraction failed: %s", exc)
            return None

    def _is_readable_text(self, text: str) -> bool:
        """Check if extracted text is readable (not garbage/binary)."""
        # Simple heuristic: check ratio of printable characters
        if not text:
            return False

        printable = sum(1 for c in text if c.isprintable() or c.isspace())
        ratio = printable / len(text)

        # Also check for reasonable word-like patterns
        words = text.split()
        avg_word_len = sum(len(w) for w in words) / max(len(words), 1)

        return ratio > 0.9 and 2 < avg_word_len < 15

    def _process_office_doc(self, record: FileRecord) -> bool:
        """Process Office document by extracting text."""
        ext = (record.file_type or "").lower()
        text = None

        try:
            if ext in {'.docx', '.doc'}:
                text = self._extract_docx_text(record.path)
            elif ext in {'.xlsx', '.xls'}:
                text = self._extract_xlsx_text(record.path)
            elif ext in {'.pptx', '.ppt'}:
                text = self._extract_pptx_text(record.path)

            if text:
                summary = self.engine._generate_text_summary(text, file_path=record.path)
                if summary:
                    from datetime import datetime

                    from duplicleaner.db.models import AISummary

                    ai_summary = AISummary(
                        file_id=record.id,
                        summary=summary,
                        document_summary=summary,
                        document_type=ext[1:],  # Remove leading dot
                        summary_model=self.config.ai.get_summary_model(),
                        generated_at=datetime.now(),
                        user_edited=False,
                    )
                    self.db.add_ai_summary(ai_summary)
                    return True

            return False

        except Exception as exc:
            logger.warning("Office doc processing failed for %s: %s", record.path, exc)
            return False

    def _extract_docx_text(self, docx_path: str) -> str | None:
        """Extract text from .docx file."""
        try:
            from docx import Document

            doc = Document(docx_path)
            text_parts = [para.text for para in doc.paragraphs]
            return "\n".join(text_parts)

        except ImportError:
            logger.warning("python-docx not installed - cannot extract .docx text")
            return None
        except Exception as exc:
            logger.warning("DOCX text extraction failed: %s", exc)
            return None

    def _extract_xlsx_text(self, xlsx_path: str) -> str | None:
        """Extract text from .xlsx file."""
        try:
            from openpyxl import load_workbook

            wb = load_workbook(xlsx_path, read_only=True, data_only=True)
            text_parts = []

            # Read first 5 sheets max
            for sheet_name in list(wb.sheetnames)[:5]:
                sheet = wb[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")

                # Read first 100 rows max per sheet
                for row in list(sheet.iter_rows(max_row=100, values_only=True)):
                    row_text = " | ".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        text_parts.append(row_text)

            wb.close()
            return "\n".join(text_parts)

        except ImportError:
            logger.warning("openpyxl not installed - cannot extract .xlsx text")
            return None
        except Exception as exc:
            logger.warning("XLSX text extraction failed: %s", exc)
            return None

    def _extract_pptx_text(self, pptx_path: str) -> str | None:
        """Extract text from .pptx file."""
        try:
            from pptx import Presentation

            prs = Presentation(pptx_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"Slide {slide_num}:")

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)

            return "\n".join(text_parts)

        except ImportError:
            logger.warning("python-pptx not installed - cannot extract .pptx text")
            return None
        except Exception as exc:
            logger.warning("PPTX text extraction failed: %s", exc)
            return None

    def _process_video_file(self, record: FileRecord) -> bool:
        """Process video by extracting keyframes and analyzing with vision model."""
        try:
            frames = self._extract_video_frames(
                record.path,
                interval=self.config.ai.video_frame_interval,
                max_frames=5,
            )
            if not frames:
                logger.warning("No frames extracted from video: %s", record.path)
                return False

            logger.info("Extracted %d frames from video: %s", len(frames), record.path)

            # Summarize each frame via the vision model
            frame_descriptions: list[str] = []
            for i, (b64, mime) in enumerate(frames):
                summary = self.engine.summarize_video_frame(b64, mime, i + 1, len(frames), file_path=record.path)
                if summary:
                    frame_descriptions.append(f"Frame {i + 1}: {summary}")

            if not frame_descriptions:
                logger.warning("Vision model returned no descriptions for: %s", record.path)
                return False

            # Combine frame descriptions into a cohesive video summary
            combined_text = "\n".join(frame_descriptions)
            video_summary = self.engine.summarize_video_combined(combined_text, file_path=record.path)
            if not video_summary:
                # Fall back to concatenated frame descriptions
                video_summary = " ".join(
                    desc.split(": ", 1)[1] if ": " in desc else desc
                    for desc in frame_descriptions
                )

            from datetime import datetime

            from duplicleaner.db.models import AISummary

            ai_summary = AISummary(
                file_id=record.id,
                summary=video_summary,
                document_type="video",
                summary_model=self.config.ai.get_summary_model(),
                generated_at=datetime.now(),
                user_edited=False,
            )
            self.db.add_ai_summary(ai_summary)
            return True

        except Exception as exc:
            logger.warning("Video processing failed for %s: %s", record.path, exc)
            return False

    def _extract_video_frames(
        self,
        video_path: str,
        interval: int = 30,
        max_frames: int = 5,
    ) -> list[tuple[str, str]]:
        """Extract frames from a video at regular intervals using OpenCV.

        Returns list of (base64_str, mime_type) tuples.
        """
        try:
            import cv2
        except ImportError:
            logger.warning("opencv-python-headless not installed - cannot extract video frames")
            return []

        import io

        from PIL import Image

        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            logger.warning("Could not open video: %s", video_path)
            return []

        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))

        if fps <= 0 or total_frames <= 0:
            logger.warning("Invalid video metadata (fps=%.1f, frames=%d): %s", fps, total_frames, video_path)
            cap.release()
            return []

        duration_sec = total_frames / fps
        frame_interval = max(interval, 1)

        # Calculate which timestamps to extract
        timestamps = []
        t = 0.0
        while t < duration_sec and len(timestamps) < max_frames:
            timestamps.append(t)
            t += frame_interval

        # Always include the first frame; if video is short, grab at least one
        if not timestamps:
            timestamps = [0.0]

        max_dim = self.config.ai.max_image_dimension
        results: list[tuple[str, str]] = []

        for ts in timestamps:
            frame_number = int(ts * fps)
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_number)
            ret, frame = cap.read()
            if not ret:
                continue

            # Convert BGR (OpenCV) -> RGB (Pillow)
            rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            img = Image.fromarray(rgb)

            # Resize to fit within max_image_dimension
            if img.width > max_dim or img.height > max_dim:
                img.thumbnail((max_dim, max_dim), Image.LANCZOS)

            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=85)
            b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
            results.append((b64, "image/jpeg"))

        cap.release()
        logger.debug(
            "Extracted %d/%d frames from %.1fs video: %s",
            len(results), len(timestamps), duration_sec, video_path,
        )
        return results

    def _process_audio_file(self, record: FileRecord) -> bool:
        """
        Process audio by transcribing with Whisper and summarizing with text model.

        """
        try:
            text = self._transcribe_audio(record.path)
            if not text:
                return False

            summary = self.engine._generate_text_summary(text, file_path=record.path)
            if not summary:
                return False

            from datetime import datetime

            from duplicleaner.db.models import AISummary

            ai_summary = AISummary(
                file_id=record.id,
                summary=summary,
                document_summary=summary,
                document_type="audio",
                summary_model=self.config.ai.get_summary_model(),
                generated_at=datetime.now(),
                user_edited=False,
            )
            self.db.add_ai_summary(ai_summary)
            return True

        except Exception as exc:
            logger.warning("Audio processing failed for %s: %s", record.path, exc)
            return False

    def _render_pdf_pages(
        self,
        pdf_path: str,
        output_dir: str,
        max_pages: int = 3,
    ) -> list[str]:
        """Render PDF pages to images and return image paths."""
        try:
            import fitz  # PyMuPDF
            from PIL import Image
        except ImportError:
            logger.warning("PyMuPDF or Pillow not installed - cannot render PDF pages")
            return []

        image_paths: list[str] = []
        doc = fitz.open(pdf_path)
        try:
            page_count = min(max_pages, len(doc))
            for page_num in range(page_count):
                page = doc[page_num]
                matrix = fitz.Matrix(2.0, 2.0)
                pix = page.get_pixmap(matrix=matrix, alpha=False)

                mode = "RGB"
                if pix.n == 1:
                    mode = "L"

                img = Image.frombytes(mode, (pix.width, pix.height), pix.samples)

                max_dim = self.config.ai.max_image_dimension
                if img.width > max_dim or img.height > max_dim:
                    img.thumbnail((max_dim, max_dim), Image.LANCZOS)

                out_path = Path(output_dir) / f"page_{page_num + 1}.jpg"
                img.save(out_path, format="JPEG", quality=85)
                image_paths.append(str(out_path))
        finally:
            doc.close()

        return image_paths

    def _combine_page_summaries(
        self, summaries: list[str], label: str, *, file_path: str = "",
    ) -> str | None:
        """Combine page summaries into a single document summary."""
        if not summaries:
            return None
        prompt = (
            f"Below are page summaries from a {label} document. "
            "Write a concise 2-3 sentence overall summary that captures the main content, "
            "key topics, and any important entities or details.\n\n"
        )
        combined_text = "\n".join(summaries)
        return self.engine._generate_text_summary(prompt + combined_text, file_path=file_path)

    def _transcribe_audio(self, audio_path: str) -> str | None:
        """Transcribe audio using Whisper if available."""
        # Prefer faster-whisper if installed
        try:
            from faster_whisper import WhisperModel

            if self._faster_whisper_model is None:
                self._faster_whisper_model = WhisperModel(
                    self.config.ai.audio_whisper_model,
                    device=self.config.ai.audio_whisper_device,
                    compute_type=self.config.ai.audio_whisper_compute_type,
                )

            segments, _info = self._faster_whisper_model.transcribe(audio_path)
            text = " ".join(segment.text for segment in segments).strip()
            return text or None
        except ImportError:
            pass
        except Exception as exc:
            logger.warning("faster-whisper transcription failed: %s", exc)

        try:
            import whisper
        except ImportError:
            logger.warning("openai-whisper not installed - cannot transcribe audio")
            return None

        if self._whisper_model is None:
            try:
                self._whisper_model = whisper.load_model(self.config.ai.audio_whisper_model)
            except Exception as exc:
                logger.warning("Failed to load Whisper model: %s", exc)
                return None

        try:
            result = self._whisper_model.transcribe(audio_path)
            text = result.get("text", "")
            return text.strip() or None
        except Exception as exc:
            logger.warning("Whisper transcription failed: %s", exc)
            return None
